#!/usr/bin/env python3
"""
Build a defense presentation from the SDSU white-background template,
report/main.tex narrative, Results-SS screenshots (authoritative numbers),
and the BERTopic modularity SVG rasterized for one slide.

Usage (from repo root):
  .venv/bin/python scripts/build_defense_slides.py
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "sdsu-white-background-powerpoint-updated.pptx"
OUT = ROOT / "Defense_Presentation_BERTopic_SDSU.pptx"
RESULTS_SS = ROOT / "Results-SS"
MODULARITY_SVG = ROOT / "report" / "figures" / "modularity.svg"


def delete_slide(prs: Presentation, slide) -> None:
    slide_id = slide.slide_id
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        if sld_id.id == slide_id:
            prs.part.drop_rel(sld_id.rId)
            sld_id_lst.remove(sld_id)
            break


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        delete_slide(prs, prs.slides[0])


def find_ss(pattern: str) -> Path:
    """Match one PNG in Results-SS; pattern is a substring of the filename."""
    matches = [p for p in RESULTS_SS.glob("*.png") if pattern in p.name]
    if not matches:
        raise FileNotFoundError(f"No screenshot matching {pattern!r} in {RESULTS_SS}")
    return sorted(matches, key=lambda p: p.name)[0]


def set_body(tf, text: str, size_pt: float = 18) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    for line in text.split("\n")[1:]:
        r = tf.add_paragraph()
        r.text = line
        r.level = 0
        r.font.size = Pt(size_pt)


def add_bullets(tf, lines: list[str], size_pt: float = 17) -> None:
    tf.clear()
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size_pt)


def add_title_only_slide(prs: Presentation, layout_idx: int, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title


def add_title_content(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    add_bullets(body, bullets, 16)


def add_picture_bottom_caption(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str,
) -> None:
    """Title on top; image fills mid; short caption at bottom (blank layout)."""
    blank = prs.slide_layouts[2]
    slide = prs.slides.add_slide(blank)
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    # Image
    slide.shapes.add_picture(
        str(image_path),
        Inches(0.55),
        Inches(1.25),
        width=Inches(12.35),
        height=Inches(5.5),
    )
    cap = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.55))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(12)
    cp.font.italic = True


def add_two_images(
    prs: Presentation,
    title: str,
    path_top: Path,
    path_bottom: Path,
    cap_top: str,
    cap_bottom: str,
) -> None:
    blank = prs.slide_layouts[2]
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.75))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    slide.shapes.add_picture(str(path_top), Inches(0.5), Inches(1.05), width=Inches(12.4), height=Inches(2.85))
    ct = slide.shapes.add_textbox(Inches(0.6), Inches(3.88), Inches(12.0), Inches(0.4))
    ct.text_frame.paragraphs[0].text = cap_top
    ct.text_frame.paragraphs[0].font.size = Pt(11)
    slide.shapes.add_picture(str(path_bottom), Inches(0.5), Inches(4.25), width=Inches(12.4), height=Inches(2.85))
    cb = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(12.0), Inches(0.4))
    cb.text_frame.paragraphs[0].text = cap_bottom
    cb.text_frame.paragraphs[0].font.size = Pt(11)


def modularity_png_path() -> Path:
    import cairosvg

    out = Path(tempfile.gettempdir()) / "bertopic_modularity_for_pptx.png"
    cairosvg.svg2png(url=str(MODULARITY_SVG), write_to=str(out), scale=2.0)
    return out


def main() -> None:
    if not TEMPLATE.is_file():
        raise SystemExit(f"Missing template: {TEMPLATE}")
    if not RESULTS_SS.is_dir():
        raise SystemExit(f"Missing folder: {RESULTS_SS}")

    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    # Layout indices: 0 Title, 1 Title and Content, 2 Blank, 3 Picture with Caption
    title_layout = prs.slide_layouts[0]

    # --- Slide 1: Title ---
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = (
        "Batch-Recurring Topic Modeling for Evolving Text Streams\n"
        "Human-in-the-Loop Monitoring and Drift Detection"
    )
    if len(s0.placeholders) > 1:
        sub = s0.placeholders[1]
        if sub.has_text_frame:
            sub.text_frame.paragraphs[0].text = (
                "Krishna Gajera · College of Engineering · San Diego State University\n"
                "Advisor: Dr. Baris Aksanli · kgajera1820@sdsu.edu"
            )
            sub.text_frame.paragraphs[0].font.size = Pt(18)

    # --- Slide 2: Motivation ---
    add_title_content(
        prs,
        "Motivation",
        [
            "Offline topic models assume a static corpus; real systems ingest continuous text.",
            "Need: repeated windows, preserved structure, drift visibility, analyst control.",
            "This work: BERTopic as the operational model; LDA and NMF as shared-metric baselines.",
        ],
    )

    # --- Slide 3: Contributions ---
    add_title_content(
        prs,
        "Contributions",
        [
            "End-to-end architecture: ingestion → modeling → monitoring → interaction.",
            "Train-and-merge BERTopic updates for temporal adaptation with versioning.",
            "Three-way evaluation (BERTopic, LDA, NMF) on coherence, diversity, silhouette.",
            "Multi-signal drift: prevalence, centroid shift, keyword JSD, topic birth/death.",
            "Governed HITL: relabel, merge, audit log, rollback; optional Ollama label polish.",
        ],
    )

    # --- Slide 4: Train-and-merge & roles ---
    add_title_content(
        prs,
        "Temporal Update: Train-and-Merge",
        [
            "Per window: train fresh BERTopic F_t on batch B_t; merge with deployed M_{t-1} using threshold τ.",
            "Merged model M_t is versioned; prior state archived for rollback.",
            "LDA/NMF train on the cumulative corpus for fair baseline metrics (topic count harmonized).",
            "BERTopic drives inference, dashboard, and HITL; baselines are comparative only.",
        ],
    )

    # --- Slide 5: Architecture ---
    add_title_content(
        prs,
        "Layered Architecture (Prefect · MLflow · FastAPI · Streamlit)",
        [
            "Orchestration: DAG tasks, batch windows, persistence, MLflow logging.",
            "Modeling: BERTopic train/merge; optional LDA/NMF; assignments + metadata.",
            "Drift: TVD on prevalence, cosine centroid shift, keyword JSD, lifecycle sets.",
            "HITL: relabel, merge, version history, rollback with audit trail.",
            "Serving: REST API; Interaction: Streamlit pages (dashboard, drift, HITL, benchmarking).",
        ],
    )

    # --- Slide 6: Orchestration ---
    add_title_content(
        prs,
        "Orchestration Pipeline (per batch)",
        [
            "Clock/data trigger → ingest & preprocess window [t_start, t_end) with batch id.",
            "Train BERTopic (or seed); merge with prior; archive baseline checkpoints.",
            "Run drift vs. previous deployed model; write alerts + metrics; expose via API.",
            "Optional: Ollama (Gemma 3 1B) for topic label rewrites (not included in timing plots).",
        ],
    )

    # --- Slide 7: BERTopic stack (modularity figure) ---
    mod_png = modularity_png_path()
    add_picture_bottom_caption(
        prs,
        "BERTopic Modularity (stack used: MiniLM → UMAP → HDBSCAN → c-TF–IDF)",
        mod_png,
        "Source: BERTopic documentation-style modularity diagram (report/figures/modularity.svg).",
    )

    # --- Slide 8: Drift ---
    add_title_content(
        prs,
        "Drift Detection & Alerts",
        [
            "Prevalence TVD between consecutive topic distributions.",
            "Per-topic centroid shift in embedding space (1 − cosine similarity).",
            "Keyword JSD on c-TF–IDF weights; topic birth/death sets.",
            "Severity tiers drive Streamlit “Drift Alerts” and MLflow artifacts.",
        ],
    )

    # --- Slide 9: HITL ---
    add_picture_bottom_caption(
        prs,
        "Human-in-the-Loop: Merge Similar Topics",
        find_ss("6.11.29"),
        "Jaccard similarity on top words; example pair T20 vs T84 (iOS) at similarity 0.25 with 21 + 14 docs.",
    )

    # --- Slide 10: Serving + inference ---
    add_title_content(
        prs,
        "Serving & Inference",
        [
            "FastAPI: catalog, trends, drift, inference, HITL mutations, metric histories.",
            "Inference: preprocess → Sentence-BERT all-MiniLM-L6-v2 (384-d) → cosine to nearest centroid.",
            "Confidence bands: high >70%, medium 40–70%, low <40% (dashboard copy).",
            "Streamlit pages: Dashboard, Topic Drilldown, Drift Alerts, HITL Editor, Inference, Benchmarking.",
        ],
    )

    # --- Slide 11: Dashboard — charts + summary strip (combined) ---
    add_two_images(
        prs,
        "Operational Dashboard: Volume & Summary",
        find_ss("6.08.22"),
        find_ss("6.08.13"),
        "Cumulative ≈1,296 docs; ~260–270 docs/batch (latest 229); topics rise to ~80.",
        "Totals: 1,296 docs, 5 batches; latest window 77 topics / 229 docs (2017-10-01 01:45–01:50).",
    )

    # --- Slide 12: Topic distribution + trends (combined) ---
    add_two_images(
        prs,
        "Topic Distribution & Temporal Topic Trends",
        find_ss("6.09.18"),
        find_ss("6.09.42"),
        "Bubble chart: long-tail sizes (e.g., Topic 45 “Halloween Party”, 4 docs).",
        "Line chart: selected topics across five batches (flight, delivery, network, etc.).",
    )

    # --- Slide 13: Average metrics + per-batch table (combined) ---
    add_two_images(
        prs,
        "Benchmarking: Averages & Per-Batch Table",
        find_ss("6.12.13"),
        find_ss("6.13.18"),
        "Means — C_v: 0.826 / 0.640 / 0.402. Diversity: 0.911 / 0.272 / 0.421. Silhouette: 0.035 / −0.016 / 0.094.",
        "Row averages — BERTopic 0.8261, 0.9111, 0.0349 · LDA 0.6400, 0.2722, −0.0165 · NMF 0.4023, 0.4214, 0.0942.",
    )

    # --- Slide 14: Temporal evaluation + training time ---
    add_two_images(
        prs,
        "Temporal Metrics & Training Time",
        find_ss("6.13.06"),
        find_ss("6.12.59"),
        "Coherence, diversity, silhouette over time; BERTopic ~11–16.5 s/batch vs LDA ~1.4 s, NMF ~1.6 s (means n=5).",
        "Mean training: BERTopic ≫ baselines (excl. Ollama); split-axis bar charts in dashboard.",
    )

    # --- Slide 15: Inference UI + examples ---
    add_two_images(
        prs,
        "Inference Page & Example Documents",
        find_ss("6.11.59"),
        find_ss("6.10.50"),
        "MiniLM 384-d embeddings + cosine to centroids; confidence bands from dashboard.",
        "Examples: 100% confidence on sample tweets with batch/time metadata.",
    )

    # --- Slide 16: Wrap-up ---
    add_title_content(
        prs,
        "Strengths, Limits, Conclusion",
        [
            "Strengths: semantic topics, governed train-and-merge, multi-signal drift, three-way metrics, API + Streamlit.",
            "Limits: batch (not streaming) updates; weak silhouette on short text; automated split is partial.",
            "Trade-off: BERTopic highest coherence/diversity but much higher fit time than LDA/NMF.",
            "Summary: reproducible pipeline (Prefect, MLflow) with screenshot-backed evidence on Twitter batches.",
            "Thank you — questions?",
        ],
    )

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
